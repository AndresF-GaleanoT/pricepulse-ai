from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK = RGBColor(0x0F, 0x0F, 0x23)
DARK2 = RGBColor(0x1A, 0x1A, 0x2E)
CARD = RGBColor(0x25, 0x25, 0x3D)
CYAN = RGBColor(0x00, 0xD2, 0xFF)
PINK = RGBColor(0xFF, 0x00, 0x88)
GREEN = RGBColor(0x00, 0xE6, 0x76)
ORANGE = RGBColor(0xFF, 0x98, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
LIGHT = RGBColor(0xCC, 0xCC, 0xDD)

def bg(slide, color=DARK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, color, radius=True):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def txt(slide, l, t, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tf

def card(slide, l, t, w, h, title, body, title_color=CYAN, title_size=20, body_size=14, bg_color=CARD):
    s = rect(slide, l, t, w, h, bg_color)
    tf = s.text_frame
    tf.word_wrap = True
    for m in [16, 16, 12, 12]:
        tf.margin_left = Pt(m)
        tf.margin_right = Pt(m)
        tf.margin_top = Pt(m)
        tf.margin_bottom = Pt(m)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = "Calibri"
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(body_size)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"
    p2.space_before = Pt(8)

def bar(slide, l, t, w, color=CYAN):
    rect(slide, l, t, w, Inches(0.04), color)

def section_title(slide, title, subtitle=None):
    txt(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7), title, size=38, color=WHITE, bold=True)
    bar(slide, Inches(0.8), Inches(1.1), Inches(2.5), CYAN)
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.3), Inches(11), Inches(0.4), subtitle, size=16, color=GRAY)

def number_card(slide, l, t, num, label, color=CYAN):
    c = rect(slide, l, t, Inches(2.5), Inches(1.8), CARD)
    tf = c.text_frame
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(12)
    tf.margin_bottom = Pt(12)
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(36)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(6)

# ===== SLIDE 1: PORTADA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK2)
rect(s, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), CYAN)
rect(s, Inches(0), Inches(3.06), Inches(13.333), Inches(0.02), PINK)
txt(s, Inches(1), Inches(1.2), Inches(11), Inches(1.2), "PRICEPULSE AI", size=60, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.4), Inches(11), Inches(0.8), "Monitoreo inteligente de precios e-commerce con 3 agentes IA", size=24, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.4), Inches(11), Inches(0.5), "FastAPI  ·  CrewAI  ·  NVIDIA NIM  ·  n8n  ·  PostgreSQL  ·  Streamlit", size=16, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(5.6), Inches(11), Inches(0.4), "github.com/AndresF-GaleanoT/pricepulse-ai", size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ===== SLIDE 2: PROBLEMA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "El Problema")
card(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(2.5), "Seguimiento manual",
     "Miles de productos cambian de precio\ncada hora. Hacer tracking manual de\nAmazon, eBay y Mercado Libre es\nimposible a escala.", title_color=ORANGE)
card(s, Inches(6.8), Inches(2.0), Inches(5.5), Inches(2.5), "Ofertas perdidas",
     "Las mejores oportunidades duran\nminutos. Sin automatizacion es\nimposible detectarlas a tiempo y\nactuar.", title_color=PINK)
card(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(1.5), "Solucion: PricePulse AI",
     "Automatiza la busqueda, analisis y reporte de precios en 3 plataformas simultaneamente. Resultado en menos de 10 segundos.",
     title_color=GREEN, title_size=22, body_size=16, bg_color=DARK2)

# ===== SLIDE 3: STACK =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Stack Tecnologico")
techs = [
    ("FastAPI", "API async", CYAN),
    ("CrewAI", "3 agentes IA", PINK),
    ("NVIDIA NIM", "Llama 3.1 70B", GREEN),
    ("SerpAPI", "Precios en vivo", ORANGE),
    ("n8n", "Orquestador", CYAN),
    ("PostgreSQL", "Historial", GREEN),
    ("Streamlit", "Dashboard", PINK),
    ("Docker", "Portable", ORANGE),
]
for i, (name, desc, color) in enumerate(techs):
    col = i % 4
    row = i // 4
    l = Inches(0.8 + col * 3.1)
    t = Inches(2.0 + row * 2.5)
    c = rect(s, l, t, Inches(2.8), Inches(2.0), CARD)
    tf = c.text_frame
    tf.margin_left = Pt(16)
    tf.margin_right = Pt(16)
    tf.margin_top = Pt(20)
    tf.margin_bottom = Pt(12)
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(26)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(14)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

# ===== SLIDE 4: ARQUITECTURA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Arquitectura del Sistema")

# Flow boxes
flow = [
    (Inches(0.3), Inches(2.0), Inches(3.0), "n8n Schedule", "Dispara a las 00, 08, 16 UTC", CYAN, "POST"),
    (Inches(3.8), Inches(2.0), Inches(3.0), "SerpAPI", "Busca en Amazon, eBay, ML", PINK, "Precios"),
    (Inches(7.3), Inches(2.0), Inches(3.0), "CrewAI Agents", "Captador > Organizador > Redactor", GREEN, "Analisis"),
    (Inches(10.8), Inches(2.0), Inches(2.2), "PostgreSQL", "Historial persistente", ORANGE, None),
]
for l, t, w, title, desc, color, arrow in flow:
    c = rect(s, l, t, w, Inches(1.8), CARD)
    tf = c.text_frame
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(16)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = GRAY
    p2.font.name = "Calibri"
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(8)
    if arrow:
        txt(s, l + Inches(0.8), t - Inches(0.5), w, Inches(0.4), arrow, size=12, color=color, align=PP_ALIGN.CENTER)

# Output boxes
out = [
    (Inches(0.3), Inches(4.5), Inches(2.5), "Google Sheets", "Backup"),
    (Inches(3.5), Inches(4.5), Inches(2.5), "Excel .xlsx", "Exportacion"),
    (Inches(6.7), Inches(4.5), Inches(2.5), "Dashboard", "Streamlit"),
    (Inches(9.9), Inches(4.5), Inches(2.5), "PDF Report", "Descargable"),
]
for l, t, w, title, desc in out:
    c = rect(s, l, t, w, Inches(1.2), DARK2)
    tf = c.text_frame
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(8)
    tf.margin_bottom = Pt(8)
    p = tf.paragraphs[0]
    p.text = f"{title}  -  {desc}"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

# Arrows between flow and output
txt(s, Inches(4.8), Inches(3.9), Inches(2), Inches(0.3), "Resultados >", size=12, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(5.9), Inches(11), Inches(0.3), "Todo corre via Docker Compose en Oracle Linux con Cloudflare Tunnel para HTTPS", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 5: FLUJO DE PRECIOS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Flujo de Precios", "Cada ejecucion busca hasta 20 precios en 3 plataformas")
cards = [
    (Inches(0.8), Inches(2.0), "Amazon", "Motor dedicado engine=amazon\n5 resultados directos\nPrecio + link limpio", CYAN),
    (Inches(4.8), Inches(2.0), "eBay", "Google Shopping filtrado\n10 resultados\nSolo source=ebay", PINK),
    (Inches(8.8), Inches(2.0), "Mercado Libre", "Google site:mercadolibre.com\n5 resultados\nPrecio via rich_snippet", GREEN),
]
for l, t, title, body, color in cards:
    card(s, l, t, Inches(3.5), Inches(2.8), title, body, title_color=color)
txt(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.4), "Cache configurable (TTL 4h)  ·  Los 10 mejores precios van a CrewAI  ·  Fallback silencioso por plataforma", size=14, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(5.5), Inches(11), Inches(0.8), "Ejemplo: Raspberry Pi 5", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.5), "Amazon: 16 resultados desde $86  |  eBay: 8 resultados desde $97  |  ML: (en desarrollo)", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 6: 3 AGENTES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Pipeline de 3 Agentes CrewAI")
# Agent cards
agents = [
    (Inches(0.8), Inches(2.0), "1  CAPTADOR", "Extrae metricas numericas:\n* Precio minimo\n* Precio maximo\n* Precio promedio\n* Anomalias detectadas\n* Mejor plataforma", CYAN),
    (Inches(4.8), Inches(2.0), "2  ORGANIZADOR", "Estructura en 4 secciones:\n* Metricas\n* Comparativa\n* Anomalias\n* Recomendacion\n* Veredicto", PINK),
    (Inches(8.8), Inches(2.0), "3  REDACTOR", "Genera JSON estricto:\n* 9 campos exactos\n* Sin markdown\n* Sin texto extra\n* Listo para API", GREEN),
]
for l, t, title, body, color in agents:
    card(s, l, t, Inches(3.5), Inches(3.5), title, body, title_color=color, title_size=22)
# Flow arrows
for x in [Inches(4.3), Inches(8.3)]:
    txt(s, x, Inches(3.2), Inches(0.5), Inches(0.4), ">", size=28, color=GRAY, align=PP_ALIGN.CENTER)
txt(s, Inches(0.8), Inches(6.0), Inches(11), Inches(0.4), "LLM: NVIDIA NIM Llama 3.1 70B  |  Temperatura: 0  |  Verbose: false  |  Timeout: 120s", size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 7: RESPUESTA JSON =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Respuesta de la API", "Estructura JSON completa para integracion automatica")
# Code box
code_box = rect(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5), DARK2)
tf = code_box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(20)
tf.margin_right = Pt(20)
tf.margin_top = Pt(16)
tf.margin_bottom = Pt(16)
code = """{
  "status": "success",
  "producto": "Raspberry Pi 5",
  "precios_encontrados": [ { "plataforma": "Amazon", "titulo": "...", "precio": 86.18, "link": "..." } ],
  "reporte_ia": "Analisis completo del Captador, Organizador y Redactor...",
  "resumen": {
    "precio_minimo": 86.18, "precio_maximo": 259.95, "precio_promedio": 169.99,
    "mejor_plataforma": "Amazon", "veredicto": "oferta", "recomendacion": "comprar"
  },
  "filas": [ { "fecha": "2026-07-30", "producto": "Raspberry Pi 5", "plataforma": "Amazon", "precio": 86.18 } ]
}"""
p = tf.paragraphs[0]
p.text = code
p.font.size = Pt(12)
p.font.color.rgb = GREEN
p.font.name = "Consolas"
card(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), "filas array -> Directo a Google Sheets y Excel via n8n",
     "Cada fila incluye: fecha, producto, plataforma, titulo, precio, link + fila RESUMEN con veredicto IA",
     title_color=ORANGE)

# ===== SLIDE 8: DASHBOARD =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Dashboard en Vivo", "Streamlit - sin dependencias externas (lee directo de la API)")
# Metrics
number_card(s, Inches(0.5), Inches(2.0), "35+", "Precios por ejecucion", CYAN)
number_card(s, Inches(3.3), Inches(2.0), "3", "Plataformas", PINK)
number_card(s, Inches(6.1), Inches(2.0), "10s", "Respuesta total", GREEN)
number_card(s, Inches(8.9), Inches(2.0), "$0", "Costo operativo", ORANGE)
# Features
feats = [
    (Inches(0.5), Inches(4.2), "Graficos", "Promedio por plataforma y tendencias historicas"),
    (Inches(3.3), Inches(4.2), "Filtros", "Por producto, plataforma y rango de fechas"),
    (Inches(6.1), Inches(4.2), "Tabla", "Detalle con links a cada oferta"),
    (Inches(8.9), Inches(4.2), "Metricas", "KPIs en tiempo real"),
]
for l, t, title, body in feats:
    card(s, l, t, Inches(3.0), Inches(1.5), title, body, title_size=18, body_size=12, bg_color=DARK2)

# ===== SLIDE 9: DEMO =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Demo Rapida")
# Box 1: curl
c1 = rect(s, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.0), DARK2)
tf1 = c1.text_frame
tf1.word_wrap = True
tf1.margin_left = Pt(16)
tf1.margin_right = Pt(16)
tf1.margin_top = Pt(16)
p = tf1.paragraphs[0]
p.text = "curl -X POST http://localhost:8000/analizar-precios \\\n  -H \"Content-Type: application/json\" \\\n  -d '{\"producto\": \"Raspberry Pi 5\",\n        \"plataformas\": [\"amazon\", \"ebay\", \"mercadolibre\"]}'"
p.font.size = Pt(13)
p.font.color.rgb = WHITE
p.font.name = "Consolas"
p2 = tf1.add_paragraph()
p2.text = "\n$ curl | jq '.filas | length'\n=> 24 resultados"
p2.font.size = Pt(13)
p2.font.color.rgb = GREEN
p2.font.name = "Consolas"
p2.space_before = Pt(12)
# Box 2: results
c2 = rect(s, Inches(6.8), Inches(2.0), Inches(5.5), Inches(4.0), DARK2)
tf2 = c2.text_frame
tf2.word_wrap = True
tf2.margin_left = Pt(16)
tf2.margin_right = Pt(16)
tf2.margin_top = Pt(16)
p = tf2.paragraphs[0]
p.text = "Resultados del analisis:"
p.font.size = Pt(16)
p.font.color.rgb = CYAN
p.font.bold = True
p.font.name = "Calibri"
for line in [
    "Amazon: 16 productos desde $86",
    "eBay: 8 productos desde $97",
    "Mejor precio: Amazon a $86 (Pi 5 2GB)",
    "Promedio general: ~$170",
    "Veredicto: OFERTA en Amazon",
    "",
    "-> 24 filas enviadas a Google Sheets",
]:
    pp = tf2.add_paragraph()
    pp.text = line
    pp.font.size = Pt(13)
    pp.font.color.rgb = GREEN if "desde" in line else GRAY
    pp.font.name = "Calibri"
    pp.space_before = Pt(4)
txt(s, Inches(0.8), Inches(6.3), Inches(11), Inches(0.4), "Pruebalo tu mismo: docker compose exec api curl http://localhost:8000/health", size=12, color=GRAY, align=PP_ALIGN.CENTER)

# ===== SLIDE 10: RESULTADOS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s)
section_title(s, "Resultados y Costos")
results_data = [
    (Inches(0.8), Inches(2.0), "34+ precios/ejecucion", "Amazon + eBay + Mercado Libre\nen menos de 10 segundos", CYAN),
    (Inches(4.8), Inches(2.0), "100% automatico", "Schedule n8n cada 8h\n0 intervencion manual", PINK),
    (Inches(8.8), Inches(2.0), "Costo $0/mes", "Todo en tiers gratuitos\nNVIDIA NIM + SerpAPI", GREEN),
]
for l, t, title, body, color in results_data:
    card(s, l, t, Inches(3.5), Inches(2.0), title, body, title_color=color, title_size=22)
# Cost table
cost_box = rect(s, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), DARK2)
tf = cost_box.text_frame
tf.word_wrap = True
tf.margin_left = Pt(16)
tf.margin_right = Pt(16)
tf.margin_top = Pt(16)
p = tf.paragraphs[0]
p.text = "Desglose de costos operativos"
p.font.size = Pt(18)
p.font.color.rgb = WHITE
p.font.bold = True
p.font.name = "Calibri"
costs = [
    "NVIDIA NIM (1000 credits/mes) ................. Gratuito",
    "SerpAPI (100 busquedas/mes) .................. Gratuito",
    "PostgreSQL (self-hosted) ..................... $0",
    "n8n (self-hosted) .............................. $0",
    "Cloudflare Tunnel .............................. $0",
    "Infraestructura (Oracle Linux) ................ $0",
]
for line in costs:
    pp = tf.add_paragraph()
    pp.text = line
    pp.font.size = Pt(13)
    pp.font.color.rgb = GRAY
    pp.font.name = "Consolas"
    pp.space_before = Pt(3)

# ===== SLIDE 11: CIERRE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK2)
# Diagonal accent
rect(s, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), CYAN)
rect(s, Inches(0), Inches(3.06), Inches(13.333), Inches(0.02), PINK)
txt(s, Inches(1), Inches(1.5), Inches(11), Inches(1.0), "GRACIAS", size=54, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(3.5), Inches(11), Inches(0.6), "github.com/AndresF-GaleanoT/pricepulse-ai", size=22, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, Inches(1), Inches(4.3), Inches(11), Inches(0.5), "Hecho con Python, NVIDIA NIM, CrewAI y mucho cafe", size=16, color=GRAY, align=PP_ALIGN.CENTER)
# Badges
badges = [
    ("Python 3.11", Inches(2.5), CYAN),
    ("FastAPI", Inches(4.5), GREEN),
    ("CrewAI", Inches(6.0), PINK),
    ("NVIDIA NIM", Inches(7.8), GREEN),
    ("Docker", Inches(10.0), CYAN),
]
for label, l, color in badges:
    b = rect(s, l, Inches(5.5), Inches(1.8), Inches(0.6), color)
    tf = b.text_frame
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(13)
    p.font.color.rgb = DARK
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

prs.save("PricePulse_AI_Presentation.pptx")
print("Presentacion mejorada generada: PricePulse_AI_Presentation.pptx")
